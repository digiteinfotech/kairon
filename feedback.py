from pptx import Presentation

prs = Presentation('./feedback/Mahesh - Comp Analysis 3 Points.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        print(shape.text)




import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import pandas as pd

headers = {"User-Agent": "Mozilla/5.0 (X11; Ubuntu; Linux x86_64; rv:82.0) Gecko/20100101 Firefox/82.0"}

vendors = [{"vendor": "AHA", "url": "https://www.getapp.com/collaboration-software/a/aha/"},
           {"vendor": "PlanView Enterprise", "url": "https://www.getapp.com/project-management-planning-software/a/planview/"},
           {"vendor": "Jira", "url": "https://www.getapp.com/project-management-planning-software/a/jira/"},
           {"vendor": "Service Now", "url": "https://www.getapp.com/it-management-software/a/service-now-com/"},
           {"vendor": "PlanView PPM Pro", "url": "https://www.getapp.com/project-management-planning-software/a/innotas/"},
           {"vendor": "Target Process", "url": "https://www.getapp.com/project-management-planning-software/a/targetprocess/"},
           {"vendor": "Version One", "url": "https://www.getapp.com/project-management-planning-software/a/versionone/"},
           {"vendor": "kanbanize", "url": "https://www.getapp.com/project-management-planning-software/a/kanbanize/"},
           {"vendor": "Clarizen", "url": "https://www.getapp.com/project-management-planning-software/a/clarizen/"},
           {"vendor": "Changepoint PPM", "url": "https://www.getapp.com/project-management-planning-software/a/daptiv-ppm/"},
           {"vendor": "Asana", "url": "https://www.getapp.com/collaboration-software/a/asana/"},
           {"vendor": "Atlassian", "url": "https://www.getapp.com/collaboration-software/a/atlassian-confluence/"},
           {"vendor": "Azure Cloud", "url": "https://www.getapp.com/it-management-software/a/windows-azure-platform/"},
           {"vendor": "WorkFront", "url": "https://www.getapp.com/project-management-planning-software/a/workfront/"},
           {"vendor": "Sciforma", "url": "https://www.getapp.com/project-management-planning-software/a/sciforma/"},
           {"vendor": "One2Team", "url": "https://www.getapp.com/project-management-planning-software/a/one2team/"},
           {"vendor": "OnePoint Projects", "url": "https://www.getapp.com/project-management-planning-software/a/onepoint-projects/"},
           {"vendor": "KeyedIn", "url": "https://www.getapp.com/project-management-planning-software/a/keyedinprojects/"},
           {"vendor": "WorkOtter", "url": "https://www.getapp.com/project-management-planning-software/a/workotter/"},
           {"vendor": "Cerri", "url": "https://www.getapp.com/project-management-planning-software/a/cerri-enterprise-apps/"}]
vendor_data = []
for vendor in vendors:
    feature_url = urljoin(vendor['url'], 'features/')
    response_features = requests.get(feature_url, headers=headers)
    vendor_features = []
    if response_features.status_code == 200:
        soup = BeautifulSoup(response_features.content, "html.parser")
        feature_pane = soup.find('div', attrs={"class": 'row feature-panes'})
        features = feature_pane.find('div', attrs={'class':'pane col-lg-4 col-md-4 col-sm-6 col-xs-12'})
        for feature in features.findAll('div', attrs={'class': 'checklist'}):
            if feature.find('i', attrs={'class': 'fa-check-circle'}):
                vendor_features.append(feature.text)

    integration_url = urljoin(vendor['url'], 'integrations/')
    response_integration = requests.get(integration_url, headers=headers)
    vendor_integration = []
    if response_integration.status_code == 200:
        soup = BeautifulSoup(response_integration.content, "html.parser")
        integrations = soup.find_all('div', attrs={'class': 'checklist'})
        for integration in integrations:
            if integration.find('i', attrs={'class': 'fa-check-circle'}):
                vendor_integration.append(integration.text)

    vendor_data.append({'Vendor': vendor['vendor'], 'features': vendor_features, "integrations": vendor_integration})

df = pd.DataFrame(vendor_data)
df.to_excel("./Vendor_Features_And_Integration.xlsx")

from sklearn.feature_extraction.text import CountVectorizer

cv = CountVectorizer(lowercase=False, tokenizer= lambda x: x)
cv_fit = cv.fit_transform(df['integrations'])

word_list = cv.get_feature_names()
count_list = cv_fit.toarray().sum(axis=0)
words = dict(zip(word_list,count_list))
final_count=[{'word':word, 'count': words[word]} for word in words]

features_count = pd.DataFrame(final_count)
features_count.to_excel('./integrations_count.xlsx')